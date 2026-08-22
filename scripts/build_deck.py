"""
Regenerable deck builder for the Foundry Pattern Library pack.
One slide per pattern — each with a native, editable architecture diagram
(rounded boxes + arrows), the money line, a one-line demo cue and the three
"where a homegrown factory falls short" chips.

Everything is drawn with python-pptx shapes (no images), so the customer's
architects can open the .pptx and tweak any box or arrow live.

Run:  uv run python scripts/build_deck.py
(If the .pptx is open in PowerPoint the file is locked; this writes a fallback name.)
"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

BLUE = RGBColor(0x0F, 0x6C, 0xBD)
DARK = RGBColor(0x20, 0x1F, 0x1E)
GREY = RGBColor(0x60, 0x5E, 0x5C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF3, 0xF2, 0xF1)
LINEG = RGBColor(0xB3, 0xB0, 0xAD)
RED = RGBColor(0xC5, 0x30, 0x1E)
AZ_TINT = RGBColor(0xEB, 0xF3, 0xFB)
AWS_TINT = RGBColor(0xF7, 0xEF, 0xE3)

W, H = Inches(13.333), Inches(7.5)

STYLES = {
    "primary": dict(fill=BLUE, text=WHITE),
    "dark": dict(fill=DARK, text=WHITE),
    "neutral": dict(fill=LIGHT, text=DARK, line=LINEG, lw=1.0),
    "accent": dict(fill=WHITE, text=BLUE, line=BLUE, lw=2.25),
    "ghost": dict(fill=WHITE, text=GREY, line=LINEG, lw=1.25, dash=True),
    "danger": dict(fill=RED, text=WHITE),
}


# ----------------------------------------------------------------------------
# primitives
# ----------------------------------------------------------------------------
def NB(x, y, w, h):
    return (Inches(x), Inches(y), Inches(w), Inches(h))


def _box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb


def _fill(slide, l, t, w, h, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def _run(p, text, size, color, bold=False, italic=False):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = "Segoe UI"
    return r


def _anchor(g, side):
    l, t, w, h = g
    if side == "l":
        return (l, t + h // 2)
    if side == "r":
        return (l + w, t + h // 2)
    if side == "t":
        return (l + w // 2, t)
    if side == "b":
        return (l + w // 2, t + h)
    return (l + w // 2, t + h // 2)


def node(slide, g, label, sub=None, style="neutral"):
    l, t, w, h = g
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    st = STYLES[style]
    sh.fill.solid()
    sh.fill.fore_color.rgb = st["fill"]
    if st.get("line"):
        sh.line.color.rgb = st["line"]
        sh.line.width = Pt(st.get("lw", 1.0))
        if st.get("dash"):
            ln = sh.line._get_or_add_ln()
            ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p, label, 12, st["text"], bold=True)
    if sub:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        _run(p2, sub, 8.5, st["text"], italic=True)
    return g


def edge(slide, g1, s1, g2, s2, label=None, dashed=False, color=GREY, width=1.75):
    p1 = _anchor(g1, s1)
    p2 = _anchor(g2, s2)
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, p1[0], p1[1], p2[0], p2[1])
    c.line.color.rgb = color
    c.line.width = Pt(width)
    ln = c.line._get_or_add_ln()
    if dashed:
        ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle"}))
    if label:
        mx = (p1[0] + p2[0]) // 2
        my = (p1[1] + p2[1]) // 2
        lb = _box(slide, mx - Inches(0.75), my - Inches(0.17), Inches(1.5), Inches(0.34))
        lb.fill.solid()
        lb.fill.fore_color.rgb = WHITE
        lb.line.fill.background()
        tf = lb.text_frame
        tf.word_wrap = False
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        pp = tf.paragraphs[0]
        pp.alignment = PP_ALIGN.CENTER
        _run(pp, label, 8.5, GREY, italic=True)
    return c


def panel(slide, g, title, tint):
    l, t, w, h = g
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = tint
    sh.line.color.rgb = LINEG
    sh.line.width = Pt(1.0)
    sh.shadow.inherit = False
    tb = _box(slide, l + Inches(0.14), t + Inches(0.06), w - Inches(0.28), Inches(0.34))
    _run(tb.text_frame.paragraphs[0], title, 11, GREY, bold=True)
    return g


# ----------------------------------------------------------------------------
# per-pattern diagrams (drawn inside the band y=2.2 .. 5.5)
# ----------------------------------------------------------------------------
def dia1(s):
    app = node(s, NB(2.5, 3.35, 2.2, 1.0), "Your app / agent", "no API key", "primary")
    gw = node(s, NB(5.7, 3.25, 2.7, 1.2), "Azure AI Gateway (APIM)", "validate-azure-ad-token", "dark")
    fdry = node(s, NB(9.9, 2.7, 2.9, 0.9), "Foundry model", "gpt-5.4-mini", "accent")
    bed = node(s, NB(9.9, 4.5, 2.9, 0.9), "AWS Bedrock", "existing provider", "ghost")
    edge(s, app, "r", gw, "l", "Entra ID token", color=BLUE)
    edge(s, gw, "r", fdry, "l", "managed identity")
    edge(s, gw, "r", bed, "l", dashed=True)


def dia2(s):
    a_agent = node(s, NB(1.0, 2.4, 3.5, 1.0), "A · Prompt agent", "rm-assistant-prompt (config)", "dark")
    a_feat = node(s, NB(1.0, 3.75, 3.5, 0.8), "File Search RAG + function tool", None, "neutral")
    b_agent = node(s, NB(8.85, 2.4, 3.5, 1.0), "B · Hosted agent", "rm-assistant-hosted (your code)", "dark")
    b_feat = node(s, NB(8.85, 3.75, 3.5, 0.8), "Your container · Foundry-managed compute", None, "neutral")
    fdry = node(s, NB(5.05, 2.5, 3.25, 0.95), "Foundry Agent Service", None, "primary")
    idn = node(s, NB(3.35, 4.65, 3.0, 0.82), "Entra Agent ID — governed", None, "accent")
    portal = node(s, NB(6.95, 4.65, 3.35, 0.82), "Portal: chat · logs · versions", None, "neutral")
    edge(s, a_feat, "t", a_agent, "b")
    edge(s, b_feat, "t", b_agent, "b")
    edge(s, a_agent, "r", fdry, "l")
    edge(s, b_agent, "l", fdry, "r")
    edge(s, fdry, "b", idn, "t")
    edge(s, fdry, "b", portal, "t")


def dia3(s):
    # The IQ family is four layers. Web IQ runs live here and Foundry IQ is the
    # enterprise half of the same story, so both are solid. Fabric IQ and Work IQ
    # are dashed: real parts of the family, not wired up in this demo.
    secret = node(s, NB(3.5, 4.78, 2.9, 0.68), "APIM secret: webiq-api-key", "injected inbound", "ghost")
    agent = node(s, NB(2.4, 3.3, 2.1, 0.95), "MCP client", "Foundry \u00b7 Copilot \u00b7 Bedrock", "dark")
    gw = node(s, NB(6.2, 3.3, 2.0, 0.95), "APIM MCP API", "authN \u00b7 quota", "primary")
    web = node(s, NB(8.5, 2.30, 3.9, 0.62), "Web IQ \u2014 cited live web", None, "neutral")
    fiq = node(s, NB(8.5, 3.10, 3.9, 0.62), "Foundry IQ \u2014 enterprise knowledge", None, "neutral")
    fab = node(s, NB(8.5, 3.90, 3.9, 0.62), "Fabric IQ \u2014 business data \u00b7 KPIs", None, "ghost")
    wiq = node(s, NB(8.5, 4.70, 3.9, 0.62), "Work IQ \u2014 M365 org context", None, "ghost")
    edge(s, agent, "r", gw, "l", "sub key only")
    edge(s, gw, "r", web, "l")
    edge(s, gw, "r", fiq, "l")
    edge(s, gw, "r", fab, "l", dashed=True)
    edge(s, gw, "r", wiq, "l", dashed=True)
    edge(s, secret, "t", gw, "b", dashed=True)


def dia4(s):
    reason = node(s, NB(5.4, 2.3, 2.6, 0.85), "Reason (Plan)", None, "primary")
    act = node(s, NB(8.7, 3.4, 2.6, 0.85), "Act — skill / tool", None, "primary")
    obs = node(s, NB(5.4, 4.55, 2.6, 0.85), "Observe", None, "primary")
    skills = node(s, NB(2.5, 2.95, 2.4, 0.8), "SKILL.md \u00d7 N", None, "neutral")
    engine = node(s, NB(2.5, 4.45, 2.4, 1.0), "Engine: Copilot SDK BYOM", "your Azure model + billing", "accent")
    edge(s, reason, "r", act, "t")
    edge(s, act, "b", obs, "r")
    edge(s, obs, "l", reason, "b")
    edge(s, skills, "r", act, "l", "capabilities")
    edge(s, engine, "r", obs, "l", dashed=True, color=BLUE)


def dia5(s):
    orch = node(s, NB(2.5, 3.35, 2.3, 0.95), "Orchestrator", None, "dark")
    an = node(s, NB(6.1, 2.4, 3.0, 0.85), "Portfolio Analyst", None, "primary")
    co = node(s, NB(6.1, 4.5, 3.0, 0.85), "Compliance Officer", None, "primary")
    agg = node(s, NB(10.1, 3.35, 2.5, 0.95), "Aggregate", None, "accent")
    edge(s, orch, "r", an, "l", "concurrent", color=BLUE)
    edge(s, orch, "r", co, "l")
    edge(s, an, "r", agg, "l")
    edge(s, co, "r", agg, "l")


def dia6(s):
    orc = node(s, NB(2.5, 2.3, 2.2, 0.75), "Foundry agent", "rm-assistant-traced", "dark")
    agn = node(s, NB(5.0, 2.3, 2.2, 0.75), "invoke_agent", None, "dark")
    tl = node(s, NB(7.5, 2.3, 2.2, 0.75), "Tool", None, "dark")
    md = node(s, NB(10.0, 2.3, 2.2, 0.75), "Model", None, "dark")
    edge(s, orc, "r", agn, "l")
    edge(s, agn, "r", tl, "l")
    edge(s, tl, "r", md, "l")
    otel = node(s, NB(4.8, 3.55, 4.2, 0.9), "OpenTelemetry spans", "tokens \u00b7 latency \u00b7 cost", "primary")
    s1 = node(s, NB(2.5, 4.55, 3.0, 0.75), "Foundry — Agents + Tracing", None, "accent")
    s2 = node(s, NB(5.85, 4.55, 3.0, 0.75), "Application Insights", None, "accent")
    s3 = node(s, NB(9.2, 4.55, 3.0, 0.75), "Datadog / Grafana", None, "ghost")
    edge(s, agn, "b", otel, "t")
    edge(s, tl, "b", otel, "t")
    edge(s, otel, "b", s1, "t")
    edge(s, otel, "b", s2, "t")
    edge(s, otel, "b", s3, "t", dashed=True)


def dia7(s):
    gs = node(s, NB(2.5, 3.3, 2.5, 1.0), "Golden set", "+ planted wrong row", "neutral")
    ag = node(s, NB(5.2, 3.35, 1.9, 0.9), "Agent", None, "dark")
    ev = node(s, NB(7.3, 3.2, 2.7, 1.2), "Evaluators", "groundedness \u00b7 tool-accuracy", "primary")
    gt = node(s, NB(10.2, 3.35, 2.4, 0.9), "CI gate", None, "accent")
    fail = node(s, NB(10.2, 4.6, 2.4, 0.75), "FAIL \u2192 blocks merge", None, "danger")
    fnd = node(s, NB(7.3, 5.0, 2.7, 0.75), "Foundry \u2014 Evaluations tab", None, "accent")
    edge(s, gs, "r", ag, "l")
    edge(s, ag, "r", ev, "l")
    edge(s, ev, "r", gt, "l")
    edge(s, gt, "b", fail, "t", color=RED)
    edge(s, ev, "b", fnd, "t")


def dia8(s):
    turn = node(s, NB(2.5, 2.45, 2.7, 0.85), "User turn", None, "neutral")
    doc = node(s, NB(2.5, 4.5, 2.7, 0.85), "Client doc — hidden XPIA", None, "neutral")
    shield = node(s, NB(5.6, 3.25, 3.0, 1.2), "Prompt Shields", "jailbreak + indirect (XPIA)", "primary")
    agent = node(s, NB(9.5, 2.55, 3.0, 0.85), "Agent \u2192 model", None, "dark")
    blocked = node(s, NB(9.5, 4.5, 3.0, 0.85), "BLOCKED", None, "danger")
    edge(s, turn, "r", shield, "l")
    edge(s, doc, "r", shield, "l")
    edge(s, shield, "r", agent, "l", "allow")
    edge(s, shield, "r", blocked, "l", "deny", color=RED)


def dia9(s):
    panel(s, NB(2.45, 2.3, 4.35, 3.2), "Microsoft Azure", AZ_TINT)
    panel(s, NB(8.05, 2.3, 4.85, 3.2), "AWS", AWS_TINT)
    fa = node(s, NB(2.8, 3.45, 3.05, 0.95), "Foundry Agent", None, "dark")
    mcp = node(s, NB(6.95, 3.45, 1.9, 0.95), "MCP / A2A", None, "primary")
    lt = node(s, NB(8.45, 2.8, 3.95, 0.8), "AWS tool — Lambda", None, "neutral")
    ba = node(s, NB(8.45, 4.5, 3.95, 0.8), "Amazon Bedrock agent", None, "neutral")
    edge(s, fa, "r", mcp, "l")
    edge(s, mcp, "r", lt, "l")
    edge(s, mcp, "r", ba, "l", "A2A", dashed=True)


def dia10(s):
    u = node(s, NB(2.4, 3.45, 1.9, 0.85), "RM / client", None, "neutral")
    ag = node(s, NB(4.6, 3.3, 2.8, 1.05), "Foundry agent", "MemorySearchPreviewTool", "dark")
    st = node(s, NB(9.1, 2.25, 3.3, 0.8), "Conversation", "short-term \u00b7 one session", "primary")
    lt = node(s, NB(9.1, 3.9, 3.3, 1.0), "Memory Store", "long-term \u00b7 per-user \u00b7 TTL", "accent")
    emb = node(s, NB(9.1, 5.0, 3.3, 0.5), "text-embedding-3-small", None, "ghost")
    edge(s, u, "r", ag, "l")
    edge(s, ag, "r", st, "l", "same session")
    edge(s, ag, "r", lt, "l", "across sessions")
    edge(s, lt, "b", emb, "t", dashed=True)


def dia11(s):
    app = node(s, NB(2.4, 3.45, 2.0, 0.85), "App", "Entra token", "neutral")
    gw = node(s, NB(4.8, 3.25, 2.8, 1.15), "Azure AI Gateway (APIM)", "semantic cache \u2014 paraphrases", "dark")
    mr = node(s, NB(9.3, 2.6, 3.1, 0.95), "Model Router", "cheapest capable model", "primary")
    m = node(s, NB(9.3, 4.1, 3.1, 1.0), "Foundry model", "prompt cache \u2192 cached_tokens", "accent")
    edge(s, app, "r", gw, "l")
    edge(s, gw, "r", mr, "l", "cache miss")
    edge(s, mr, "b", m, "t")


def dia12(s):
    runs = node(s, NB(2.4, 2.5, 3.0, 0.95), "Live agent runs", "tokens + outcome", "primary")
    tr = node(s, NB(2.4, 4.35, 3.0, 0.95), "App Insights traces", "Pattern 6 \u00b7 KQL by agent", "ghost")
    roi = node(s, NB(6.7, 3.25, 3.0, 1.15), "Cost \u2194 value \u2194 ROI", "per agent, per version", "dark")
    a365 = node(s, NB(10.3, 3.3, 2.3, 1.05), "Agent 365", "inventory \u00b7 identity \u00b7 policy", "accent")
    edge(s, runs, "r", roi, "l")
    edge(s, tr, "r", roi, "l", dashed=True)
    edge(s, a365, "l", roi, "r")


# ----------------------------------------------------------------------------
# slide chrome
# ----------------------------------------------------------------------------
def _chips(s, beats):
    xs = [2.45, 6.0, 9.55]
    for i, txt in enumerate(beats[:3]):
        l, t, w, h = NB(xs[i], 5.95, 3.4, 0.9)
        sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
        sh.fill.solid()
        sh.fill.fore_color.rgb = LIGHT
        sh.line.color.rgb = RGBColor(0xE1, 0xDF, 0xDD)
        sh.line.width = Pt(0.75)
        sh.shadow.inherit = False
        tf = sh.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.12)
        tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.04)
        tf.margin_bottom = Inches(0.04)
        p = tf.paragraphs[0]
        _run(p, "\u25aa  ", 10.5, BLUE, bold=True)
        _run(p, txt, 10.5, DARK)


def pattern_slide(prs, num, title, money, demo, beats, diag):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _fill(s, 0, 0, Inches(2.1), H, DARK)
    n = _box(s, Inches(0.25), Inches(0.4), Inches(1.7), Inches(1.3))
    _run(n.text_frame.paragraphs[0], f"{num:02d}", 54, BLUE, bold=True)
    tag = _box(s, Inches(0.25), Inches(1.75), Inches(1.7), Inches(0.5))
    _run(tag.text_frame.paragraphs[0], "PATTERN", 12, RGBColor(0xA1, 0x9F, 0x9D), bold=True)

    t = _box(s, Inches(2.35), Inches(0.28), Inches(10.7), Inches(0.95))
    _run(t.text_frame.paragraphs[0], title, 24, DARK, bold=True)

    m = _box(s, Inches(2.4), Inches(1.28), Inches(10.6), Inches(0.5))
    _run(m.text_frame.paragraphs[0], f"\u201c{money}\u201d", 15, BLUE, italic=True, bold=True)

    d = _box(s, Inches(2.4), Inches(1.78), Inches(10.6), Inches(0.42))
    dp = d.text_frame.paragraphs[0]
    _run(dp, "DEMO   ", 10, BLUE, bold=True)
    _run(dp, demo, 11.5, GREY)

    diag(s)

    lab = _box(s, Inches(2.45), Inches(5.64), Inches(10.4), Inches(0.28))
    _run(lab.text_frame.paragraphs[0], "WHERE A HOMEGROWN FACTORY FALLS SHORT", 9, GREY, bold=True)
    _chips(s, beats)


def group_slide(prs, idx, total, name, blurb, members):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _fill(s, 0, 0, W, H, DARK)
    _fill(s, 0, Inches(2.15), Inches(0.28), Inches(1.45), BLUE)
    tag = _box(s, Inches(0.7), Inches(1.72), Inches(11), Inches(0.4))
    _run(tag.text_frame.paragraphs[0], f"GROUP {idx} OF {total}", 12, RGBColor(0xA1, 0x9F, 0x9D), bold=True)
    tb = _box(s, Inches(0.7), Inches(2.08), Inches(12), Inches(1.3))
    _run(tb.text_frame.paragraphs[0], name, 40, WHITE, bold=True)
    p2 = tb.text_frame.add_paragraph()
    _run(p2, blurb, 20, BLUE, bold=True)
    body = _box(s, Inches(0.78), Inches(4.0), Inches(12), Inches(2.6))
    tf = body.text_frame
    first = True
    for num, title in members:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(10)
        _run(p, f"{num:02d}   ", 17, BLUE, bold=True)
        _run(p, title, 17, RGBColor(0xC8, 0xC6, 0xC4))


def title_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _fill(s, 0, 0, W, H, DARK)
    _fill(s, 0, Inches(2.55), Inches(0.28), Inches(1.7), BLUE)
    tb = _box(s, Inches(0.7), Inches(2.4), Inches(12), Inches(2))
    _run(tb.text_frame.paragraphs[0], "From AI Gateway to AI Factory", 46, WHITE, bold=True)
    p2 = tb.text_frame.add_paragraph()
    _run(p2, "Twelve Microsoft Foundry patterns — alongside AWS + your gateway", 22, RGBColor(0xC8, 0xC6, 0xC4))
    p3 = tb.text_frame.add_paragraph()
    _run(p3, "A Private Banking scenario  \u00b7  Build in GitHub \u2192 Run in Foundry \u2192 Reach in M365", 16, BLUE, bold=True)


def story_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _fill(s, 0, 0, W, Inches(1.2), BLUE)
    t = _box(s, Inches(0.6), Inches(0.28), Inches(12), Inches(0.8))
    _run(t.text_frame.paragraphs[0], "The wedge: keep your gateway, keep AWS — add the factory", 30, WHITE, bold=True)

    body = _box(s, Inches(0.7), Inches(1.5), Inches(12), Inches(3.2))
    tf = body.text_frame
    lines = [
        ("A gateway gives you MODEL ACCESS.", True),
        ("Your LiteLLM gateway (or Azure API Management) is the right pattern for that — one OpenAI-compatible endpoint, keys, budgets, routing across providers incl. Bedrock.", False),
        ("Foundry gives you the AGENT FACTORY.", True),
        ("The runtime (Plan/Act/Observe), grounding (Microsoft IQ), identity, evaluation, tracing and safety plane AROUND the models — exactly the depth a homegrown factory lacks.", False),
    ]
    first = True
    for text, bold in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(8)
        _run(p, text, 19 if bold else 15, DARK if bold else GREY, bold=bold)

    gw = node(s, NB(1.0, 5.2, 3.6, 1.1), "Your gateway", "LiteLLM / APIM \u2014 model access", "dark")
    fac = node(s, NB(5.6, 5.0, 7.0, 1.5), "Foundry \u2014 the agent factory",
               "runtime \u00b7 Microsoft IQ \u00b7 identity \u00b7 eval \u00b7 tracing \u00b7 safety", "accent")
    edge(s, gw, "r", fac, "l", "add", color=BLUE, width=2.25)


def map_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _fill(s, 0, 0, W, Inches(1.2), DARK)
    t = _box(s, Inches(0.6), Inches(0.28), Inches(12), Inches(0.8))
    _run(t.text_frame.paragraphs[0], "Run-of-show \u2014 four groups, 60 minutes", 30, WHITE, bold=True)
    rows = [
        "CONTROL PLANE \u2014 make the gateway the single front door",
        "     1  Wedge \u2192 AI Hub Gateway / Citadel  (0\u20134)   \u00b7   8  Governance / Prompt Shields  (4\u201310)",
        "AGENT FACTORY \u2014 build the agent",
        "     2  Agent Service  (10\u201315)   \u00b7   3  Microsoft IQ  (15\u201320)",
        "     4  Agentic Loop  (20\u201325)   \u00b7   10  Memory \u2014 short + long term  (25\u201329)",
        "ORCHESTRATION & INTEROP \u2014 make agents work together",
        "     5  Multi-agent (Agent Framework)  (29\u201334)   \u00b7   9  AWS cross-cloud, MCP / A2A  (34\u201340)",
        "OPERATE & OPTIMISE \u2014 run it in production",
        "     6  Observability & tracing  (40\u201347)   \u00b7   7  Evaluation \u2192 optimization  (47\u201353)",
        "     11  Caching & Cost  (53\u201357)   \u00b7   12  Agent 365 & ROI  (57\u201360)",
    ]
    body = _box(s, Inches(0.8), Inches(1.7), Inches(11.8), Inches(5.2))
    tf = body.text_frame
    first = True
    for i, r in enumerate(rows):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        is_group = not r.startswith("  ")
        # only the last member row of a group gets the wider gap before the next header
        last_of_group = (not is_group) and (i + 1 == len(rows) or not rows[i + 1].startswith("  "))
        p.space_after = Pt(14) if last_of_group else Pt(4)
        _run(p, r, 15 if is_group else 16, BLUE if is_group else DARK, bold=is_group)


def close_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _fill(s, 0, 0, W, H, DARK)
    t = _box(s, Inches(0.7), Inches(0.6), Inches(12), Inches(1))
    _run(t.text_frame.paragraphs[0], "Coexistence — and the close", 32, WHITE, bold=True)
    body = _box(s, Inches(0.8), Inches(1.8), Inches(11.8), Inches(5))
    tf = body.text_frame
    lines = [
        ("Bedrock runs an agent. Foundry runs the agent FACTORY. MCP + A2A join the two.", True),
        ("Keep Bedrock where it works \u00b7 add Foundry where you have gaps (identity, eval, tracing, grounding, safety) \u00b7 unify governance with Purview across both clouds.", False),
        ("", False),
        ("Foundry  +  Citadel  +  Agentic Patterns.", True),
        ("Platform. Governance at scale. The fastest path to business value.", False),
    ]
    first = True
    for text, bold in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(12)
        _run(p, text, 22 if bold else 16, (BLUE if bold else RGBColor(0xC8, 0xC6, 0xC4)), bold=bold)


PATTERNS = [
    (1, "The Wedge \u2192 AI Hub Gateway / Citadel",
     "The gateway gives you model access. Foundry gives you the agent factory.",
     "One OpenAI-shaped call, keyless via Entra ID \u2014 401 without a token, 200 with.",
     ["Their LiteLLM gateway stays \u2014 zero rip-and-replace",
      "Reframes gateway as commodity; factory as the differentiation",
      "Ships as AI Hub Gateway / Citadel (Foundry + APIM)"], dia1),
    (2, "Agent Service",
     "You didn't build a runtime, a vector store and an identity system. You called an API \u2014 or handed us a container.",
     "Two hosting models \u2014 declarative prompt agent + BYO-code hosted agent \u2014 both first-class, both with an Entra Agent ID.",
     ["Prompt-based: model + instructions + tools; File Search RAG (managed vector store)",
      "Hosted: your container on Foundry-managed compute \u2014 any framework",
      "Entra Agent ID \u2014 governable identity per agent, not a shared IAM role"], dia2),
    (3, "Microsoft IQ \u2014 The Intelligence Layer",
     "Web IQ grounds you in the world. Foundry IQ grounds you in your enterprise.",
     "Web IQ published as OUR OWN MCP API on APIM \u2014 no key, 401; valid key, 200; past the limit, 429.",
     ["Tool calls governed like model calls \u2014 same gateway, same control point",
      "Key custody sits in the gateway \u2014 no Web IQ credential client-side",
      "Four layers: web, enterprise, business data, org context \u2014 a moat AWS can't match"], dia3),
    (4, "Agentic Loop \u2014 Build Skills, Not Agents",
     "Don't orchestrate fifty agents. Give one good loop the right skills \u2014 and let it reason.",
     "One Plan/Act/Observe loop, N skills-as-folders; engine swapped to Copilot SDK BYOM.",
     ["One loop to maintain vs a brittle handoff graph",
      "Skills portable across tools/models; add one in minutes",
      "BYOM \u2014 proven Copilot loop, your model + your billing"], dia4),
    (5, "Multi-agent Orchestration (Agent Framework)",
     "Orchestrate specialists when the problem is genuinely parallel \u2014 not because a framework let you.",
     "Request fans out to two specialists concurrently; Compliance returns BLOCK, then aggregate.",
     ["Managed orchestration (concurrent/sequential/handoff/Magentic)",
      "Open, code-first (SK/AutoGen lineage), model-portable",
      "Honest guidance: default to Pattern 4; escalate when earned"], dia5),
    (6, "Observability & Tracing",
     "A real Foundry agent, then one traced turn \u2014 same run in the portal AND App Insights. And it's OTel, so it's yours.",
     "Run enable_tracing.py \u2014 agent 'rm-assistant-traced' + its trace in BOTH Foundry Tracing and App Insights.",
     ["Agent shows in the Agents list; run traced server-side",
      "Token + cost + latency per span \u2014 observability & FinOps",
      "OpenTelemetry \u2014 portable to Datadog/Grafana, no lock-in"], dia6),
    (7, "Evaluation \u2192 Optimization",
     "If you can't score it, you can't ship it safely \u2014 and you can't optimize it.",
     "Score the golden set locally AND in Foundry; a wrong 'suitable' tanks groundedness and the CI gate blocks the PR.",
     ["Agent-grade evaluators (groundedness, tool-call, intent)",
      "Scores locally + uploads to the Foundry Evaluations tab",
      "CI gate on every PR \u2014 regressions can't merge"], dia7),
    (8, "Governance / Prompt Shields / Content Safety",
     "Your gateway checks tokens. Foundry checks the attack \u2014 even the one hidden in a document.",
     "LIVE + keyless: Prompt Shields blocks a direct jailbreak and a document-hidden XPIA; a clean question passes.",
     ["XPIA / indirect-injection defence \u2014 rarely built in-house",
      "Entra Agent ID \u2014 scoped, conditional-access identity",
      "Purview DSPM for AI \u2014 one DLP/audit plane across clouds"], dia8),
    (9, "AWS Cross-cloud Interop (the close)",
     "Bedrock runs an agent. Foundry runs the agent factory. MCP and A2A let each cloud do what it's best at.",
     "Foundry agent calls an AWS tool over MCP (mock Lambda); A2A hands off to a Bedrock agent.",
     ["Wrap, don't rewrite \u2014 call AWS tools from Foundry via MCP",
      "A2A hand-off exposes Foundry agents to Bedrock and back",
      "Governance overlay (Entra + Purview) spans both clouds"], dia9),
    (10, "Memory \u2014 Short-term + Long-term",
     "Memory is a platform primitive \u2014 not a database you build.",
     "Recall inside one session, then a BRAND-NEW conversation honours the client's preference \u2014 nothing re-told.",
     ["Extraction, consolidation, retrieval and TTL \u2014 Foundry's job",
      "Per-user scope \u2014 isolation plus GDPR forget (delete_scope)",
      "No state store, no vector DB, no glue code to maintain"], dia10),
    (11, "Caching & Cost",
     "Two cache layers and a router \u2014 cheaper without touching your app.",
     "Repeat the stable prefix and cached_tokens jumps; a trivial prompt downshifts through Model Router.",
     ["Prompt caching \u2014 automatic; cached_tokens proves the hit",
      "Model Router \u2014 cost-aware model choice as a deployment",
      "Semantic cache at the gateway \u2014 dedupes paraphrases"], dia11),
    (12, "Agent 365 & ROI",
     "Every agent is an identity you can govern and a cost you can justify.",
     "Live RM tasks through the gateway \u2192 a cost \u2194 value \u2194 ROI table, projected to monthly volume.",
     ["Inventory + identity \u2014 Agent 365 and Entra Agent ID",
      "Cost tied to completed outcomes \u2014 a CFO-legible number",
      "Built on telemetry you already have \u2014 the Pattern 6 traces"], dia12),
]


GROUPS = [
    ("Control plane", "make the gateway the single front door", [1, 8]),
    ("Agent factory", "build the agent", [2, 3, 4, 10]),
    ("Orchestration & interop", "make agents work together", [5, 9]),
    ("Operate & optimise", "run it in production", [6, 7, 11, 12]),
]


def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    title_slide(prs)
    story_slide(prs)
    map_slide(prs)
    by_num = {p[0]: p for p in PATTERNS}
    missing = [n for _, _, nums in GROUPS for n in nums if n not in by_num]
    ungrouped = [n for n in by_num if n not in {x for _, _, nums in GROUPS for x in nums}]
    if missing or ungrouped:
        raise SystemExit(f"GROUPS/PATTERNS mismatch — missing: {missing}, ungrouped: {ungrouped}")
    for i, (name, blurb, nums) in enumerate(GROUPS, 1):
        group_slide(prs, i, len(GROUPS), name, blurb, [(n, by_num[n][1]) for n in nums])
        for n in nums:
            num, title, money, demo, beats, diag = by_num[n]
            pattern_slide(prs, num, title, money, demo, beats, diag)
    close_slide(prs)

    out = os.path.join(os.path.dirname(os.path.dirname(__file__)), "foundry-patterns.pptx")
    n = len(prs.slides._sldIdLst)
    try:
        prs.save(out)
    except PermissionError:
        out = out.replace(".pptx", ".new.pptx")
        prs.save(out)
        print("(target was locked \u2014 saved as)", end=" ")
    print(f"Saved {out}  ({n} slides)")


if __name__ == "__main__":
    build()
