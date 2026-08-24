"""Cross-check canonical pattern names, folders, group order, docs, and deck titles."""
from pathlib import Path
import re
import sys

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]

PATTERNS = [
    (1, "01-ai-gateway-model-access", "AI gateway & model access (APIM)", "Platform foundation & governance"),
    (8, "08-ai-safety", "AI safety (Prompt Shields + Content Safety)", "Platform foundation & governance"),
    (
        13,
        "13-human-approval",
        "Human approval for consequential tool actions",
        "Platform foundation & governance",
    ),
    (
        2,
        "02-foundry-agent-service",
        "Foundry Agent Service (prompt and hosted agents)",
        "Agent construction & knowledge",
    ),
    (3, "03-microsoft-iq", "Microsoft IQ \u2014 the intelligence layer", "Agent construction & knowledge"),
    (
        12,
        "12-toolbox",
        "Centralized Toolboxes (one governed MCP endpoint)",
        "Agent construction & knowledge",
    ),
    (
        14,
        "14-model-adaptation",
        "Model adaptation (fine-tuning & evaluation)",
        "Agent construction & knowledge",
    ),
    (10, "10-memory", "Memory (short-term + long-term)", "Agent construction & knowledge"),
    (
        4,
        "04-agentic-loop",
        "Agentic Loop (build skills, not agents)",
        "Orchestration & interoperability",
    ),
    (
        5,
        "05-multi-agent",
        "Multi-agent orchestration (Agent Framework)",
        "Orchestration & interoperability",
    ),
    (9, "09-aws-interop", "Cross-cloud interop (MCP / A2A)", "Orchestration & interoperability"),
    (
        7,
        "07-evaluation-release-gate",
        "Evaluation & release gate",
        "Lifecycle, assurance & operations",
    ),
    (
        6,
        "06-observability",
        "Observability & tracing (OpenTelemetry)",
        "Lifecycle, assurance & operations",
    ),
    (
        11,
        "11-caching-cost",
        "Cost & latency (prompt cache + Model Router)",
        "Lifecycle, assurance & operations",
    ),
    (
        15,
        "15-agent-lifecycle",
        "Agent lifecycle & promotion (dev \u2192 test \u2192 prod)",
        "Lifecycle, assurance & operations",
    ),
]

OLD_REFERENCES = (
    "01-wedge",
    "02-agent-service",
    "07-evaluations",
    "08-governance",
    "Wedge \u2192 AI Hub Gateway / Citadel",
    "Agent Service (prompt and hosted agent)",
    "Governance (Prompt Shields + Content Safety)",
    "Evaluation \u2192 optimization (CI gate)",
)


def fail(message: str):
    raise AssertionError(message)


def check_folders_and_talk_tracks():
    numbered = sorted(
        int(path.name[:2])
        for path in ROOT.iterdir()
        if path.is_dir() and re.fullmatch(r"\d{2}-.+", path.name)
    )
    if numbered != list(range(1, 16)):
        fail(f"numbered folders are not contiguous 1-15: {numbered}")

    for position, (number, folder, name, group) in enumerate(PATTERNS, start=1):
        talk_track = ROOT / folder / "TALK-TRACK.md"
        if not talk_track.is_file():
            fail(f"missing talk track: {talk_track.relative_to(ROOT)}")
        lines = talk_track.read_text(encoding="utf-8").splitlines()
        expected_h1 = f"# Pattern {number} \u2014 {name}"
        if lines[0] != expected_h1:
            fail(f"{folder}/TALK-TRACK.md H1: {lines[0]!r}, expected {expected_h1!r}")
        expected_group = f"**Group:** {group}"
        if not lines[2].startswith(expected_group):
            fail(f"{folder}/TALK-TRACK.md group/run header is inconsistent: {lines[2]!r}")
        if f"**Runs {ordinal(position)} of 15**" not in lines[2]:
            fail(f"{folder}/TALK-TRACK.md run position is inconsistent: {lines[2]!r}")


def ordinal(value: int) -> str:
    if 10 <= value % 100 <= 20:
        suffix = "th"
    else:
        suffix = {1: "st", 2: "nd", 3: "rd"}.get(value % 10, "th")
    return f"{value}{suffix}"


def check_readme():
    readme = (ROOT / "README.md").read_text(encoding="utf-8")
    inside = readme.split("## What's inside", 1)[1].split("### The gateway thread", 1)[0]
    diagrams = readme.split("## Pattern diagrams", 1)[1].split("## Suggested run-of-show", 1)[0]
    run_show = readme.split("## Suggested run-of-show", 1)[1]

    expected_numbers = [item[0] for item in PATTERNS]
    table_numbers = [
        int(value)
        for value in re.findall(r"^\| (\d+) \| `\d{2}-", inside, flags=re.MULTILINE)
    ]
    if table_numbers != expected_numbers:
        fail(f"README group-table order {table_numbers}, expected {expected_numbers}")
    diagram_numbers = [
        int(value)
        for value in re.findall(r"^### (\d+) \u00b7 ", diagrams, flags=re.MULTILINE)
    ]
    if diagram_numbers != expected_numbers:
        fail(f"README diagram order {diagram_numbers}, expected {expected_numbers}")
    run_numbers = [
        int(value)
        for value in re.findall(
            r"^\| [^|]+ \| (\d+) \| [^|]+ \|$",
            run_show,
            flags=re.MULTILINE,
        )
    ]
    if run_numbers != expected_numbers:
        fail(f"README run-of-show order {run_numbers}, expected {expected_numbers}")
    if "Fifteen patterns in four groups." not in readme:
        fail("README pattern count is not fifteen")

    for number, folder, name, group in PATTERNS:
        table_prefix = f"| {number} | `{folder}/` | {name} |"
        if table_prefix not in inside:
            fail(f"README group table missing: {table_prefix}")
        if f"### {number} \u00b7 {name}" not in diagrams:
            fail(f"README diagram heading missing: {number} {name}")
        run_row = f"| {group} | {number} | {name} |"
        if run_row not in run_show:
            fail(f"README run-of-show missing: {run_row}")


def slide_text(slide) -> str:
    values = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            values.append(shape.text)
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                values.extend(cell.text for cell in row.cells)
    return "\n".join(values)


def check_deck():
    deck = Presentation(ROOT / "foundry-patterns.pptx")
    if len(deck.slides) != 23:
        fail(f"deck has {len(deck.slides)} slides, expected 23")

    texts = [slide_text(slide) for slide in deck.slides]
    found = []
    for number, _, name, _ in PATTERNS:
        matches = [
            index
            for index, (slide, text) in enumerate(zip(deck.slides, texts), start=1)
            if any(
                getattr(shape, "has_text_frame", False)
                and shape.top / 914400 < 1.5
                and shape.text == f"{number:02d}"
                for shape in slide.shapes
            )
            and name in text
        ]
        if len(matches) != 1:
            fail(f"deck title {number} {name!r} found on slides {matches}")
        found.append(matches[0])
    expected_slides = [5, 6, 7, 9, 10, 11, 12, 13, 15, 16, 17, 19, 20, 21, 22]
    if found != expected_slides:
        fail(f"deck pattern slide order is inconsistent: {found}")

    all_text = "\n".join(texts)
    if "15 patterns" not in texts[0]:
        fail("deck title slide does not say 15 patterns")
    for group in dict.fromkeys(item[3] for item in PATTERNS):
        if group not in all_text:
            fail(f"deck section group missing: {group}")
    for slide_number, slide in enumerate(deck.slides, start=1):
        footer_numbers = [
            shape.text.strip()
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
            and shape.left / 914400 > 11.5
            and shape.top / 914400 > 6.8
            and shape.text.strip().isdigit()
        ]
        if footer_numbers and footer_numbers != [str(slide_number)]:
            fail(
                f"deck footer on slide {slide_number} is {footer_numbers}, "
                f"expected {slide_number}"
            )


def check_stale_references():
    extensions = {".md", ".py", ".toml", ".yml", ".yaml", ".example"}
    for path in ROOT.rglob("*"):
        if not path.is_file() or path.suffix not in extensions:
            continue
        if ".venv" in path.parts or path in {
            Path(__file__),
            ROOT / "scripts" / "refresh_deck.py",
        }:
            continue
        text = path.read_text(encoding="utf-8")
        for old in OLD_REFERENCES:
            if old in text:
                fail(f"stale reference {old!r} in {path.relative_to(ROOT)}")


def main():
    check_folders_and_talk_tracks()
    check_readme()
    check_deck()
    check_stale_references()
    print("Consistency check passed: 15 folders, talk tracks, README surfaces, and deck align.")


if __name__ == "__main__":
    try:
        main()
    except AssertionError as error:
        print(f"CONSISTENCY ERROR: {error}", file=sys.stderr)
        raise SystemExit(1) from error
