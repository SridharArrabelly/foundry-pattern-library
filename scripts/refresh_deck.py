"""Refresh the established deck without changing its purple visual style."""
from copy import deepcopy
from pathlib import Path
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
DECK = ROOT / "foundry-patterns.pptx"
PURPLE = RGBColor(0x86, 0x61, 0xC5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NAVY = RGBColor(0x09, 0x1F, 0x2C)
GRAY = RGBColor(0x60, 0x5E, 0x5C)
LIGHT_GRAY = RGBColor(0xF7, 0xF7, 0xF7)
LIGHT_PURPLE = RGBColor(0xF3, 0xF1, 0xF8)
BLUE = RGBColor(0x00, 0x78, 0xD4)
GREEN = RGBColor(0x10, 0x7C, 0x41)
ORANGE = RGBColor(0xD8, 0x9B, 0x6A)


def shape_texts(slide):
    return [shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False)]


def slide_text(slide):
    values = shape_texts(slide)
    for shape in slide.shapes:
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                values.extend(cell.text for cell in row.cells)
    return "\n".join(values)


def set_text(target, value: str):
    text_frame = target.text_frame
    source_paragraph = text_frame.paragraphs[0]
    source_run = source_paragraph.runs[0] if source_paragraph.runs else None
    run_properties = (
        deepcopy(source_run._r.get_or_add_rPr()) if source_run is not None else None
    )
    paragraph_properties = (
        deepcopy(source_paragraph._p.pPr)
        if source_paragraph._p.pPr is not None
        else None
    )

    text_frame.clear()
    for index, line in enumerate(value.split("\n")):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        if index > 0 and paragraph_properties is not None:
            existing = paragraph._p.pPr
            if existing is not None:
                paragraph._p.remove(existing)
            paragraph._p.insert(0, deepcopy(paragraph_properties))
        run = paragraph.add_run()
        run.text = line
        if run_properties is not None:
            existing = run._r.get_or_add_rPr()
            existing.getparent().replace(existing, deepcopy(run_properties))


def set_font_size(target, points):
    for paragraph in target.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(points)


def add_text(
    slide,
    text,
    left,
    top,
    width,
    height,
    *,
    size,
    color=NAVY,
    bold=False,
    align=PP_ALIGN.LEFT,
    font_name="Aptos",
    vertical_anchor=MSO_ANCHOR.TOP,
):
    shape = slide.shapes.add_textbox(
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = vertical_anchor
    for index, line in enumerate(text.split("\n")):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = align
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)
        run = paragraph.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return shape


def add_box(
    slide,
    left,
    top,
    width,
    height,
    *,
    fill=WHITE,
    line=PURPLE,
    radius=True,
    line_width=1.25,
):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(line_width)
    return shape


def clear_slide(slide):
    for shape in list(slide.shapes):
        slide.shapes._spTree.remove(shape.element)


def add_brand(slide):
    colors = (
        RGBColor(0xF2, 0x50, 0x22),
        RGBColor(0x7F, 0xBA, 0x00),
        RGBColor(0x00, 0xA4, 0xEF),
        RGBColor(0xFF, 0xB9, 0x00),
    )
    positions = ((10.88, 0.42), (11.01, 0.42), (10.88, 0.55), (11.01, 0.55))
    for color, (left, top) in zip(colors, positions):
        square = add_box(
            slide,
            left,
            top,
            0.11,
            0.11,
            fill=color,
            line=color,
            radius=False,
            line_width=0,
        )
        square.line.fill.background()
    add_text(
        slide,
        "Microsoft Foundry",
        11.24,
        0.36,
        1.55,
        0.36,
        size=11,
        color=GRAY,
        bold=True,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )


def add_footer(slide, number):
    line = add_box(
        slide,
        0.62,
        7.0,
        12.1,
        0.01,
        fill=RGBColor(0xD2, 0xD0, 0xCE),
        line=RGBColor(0xD2, 0xD0, 0xCE),
        radius=False,
        line_width=0,
    )
    line.line.fill.background()
    add_text(
        slide,
        "Microsoft Foundry Patterns  \u00b7  subject to change",
        0.62,
        7.06,
        10.5,
        0.28,
        size=9,
        color=GRAY,
        font_name="Aptos Mono",
    )
    add_text(
        slide,
        str(number),
        12.2,
        7.06,
        0.5,
        0.28,
        size=9,
        color=GRAY,
        align=PP_ALIGN.RIGHT,
        font_name="Aptos Mono",
    )


def replace_text(slide, old_values, new_value):
    if isinstance(old_values, str):
        old_values = (old_values,)
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.text in old_values:
            set_text(shape, new_value)
            return shape
    raise ValueError(f"Could not find any of {old_values!r} on slide")


def pattern_slides(presentation):
    result = {}
    for slide in presentation.slides:
        for shape in slide.shapes:
            if (
                getattr(shape, "has_text_frame", False)
                and shape.top / 914400 < 1.5
                and re.fullmatch(r"\d{2}[AB]?", shape.text)
            ):
                match = re.fullmatch(r"(\d{2})([AB]?)", shape.text)
                key = (
                    int(match.group(1))
                    if not match.group(2)
                    else f"{int(match.group(1))}{match.group(2)}"
                )
                result[key] = slide
    return result


def group_slides(presentation):
    result = {}
    for slide in presentation.slides:
        for value in shape_texts(slide):
            match = re.fullmatch(r"Group (\d) of 4", value)
            if match:
                result[int(match.group(1))] = slide
    return result


def find_slide(presentation, needle):
    return next(slide for slide in presentation.slides if needle in slide_text(slide))


def deck_label(number):
    value = str(number)
    match = re.fullmatch(r"(\d+)([AB]?)", value)
    if match is None:
        raise ValueError(f"Invalid pattern label: {number}")
    return f"{int(match.group(1)):02d}{match.group(2)}"


def clone_shape_slide(presentation, template):
    """Clone a shape-only pattern slide; the deck uses no image relationships here."""
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    for shape in list(slide.shapes):
        slide.shapes._spTree.remove(shape.element)
    for shape in template.shapes:
        slide.shapes._spTree.insert_element_before(
            deepcopy(shape.element),
            "p:extLst",
        )
    return slide


def add_stage_card(slide, index, title, body, left, *, emphasized=False):
    fill = PURPLE if emphasized else LIGHT_PURPLE
    text_color = WHITE if emphasized else NAVY
    muted = WHITE if emphasized else GRAY
    add_box(
        slide,
        left,
        2.0,
        1.75,
        2.45,
        fill=fill,
        line=PURPLE,
        line_width=1.5,
    )
    add_text(
        slide,
        f"{index:02d}",
        left + 0.18,
        2.2,
        0.45,
        0.28,
        size=10,
        color=WHITE if emphasized else PURPLE,
        bold=True,
        font_name="Aptos Mono",
    )
    add_text(
        slide,
        title,
        left + 0.18,
        2.6,
        1.39,
        0.62,
        size=15 if len(title) < 18 else 13.5,
        color=text_color,
        bold=True,
    )
    add_text(
        slide,
        body,
        left + 0.18,
        3.35,
        1.39,
        0.72,
        size=10.5,
        color=muted,
    )


def add_value_card(slide, title, body, left):
    add_box(
        slide,
        left,
        5.02,
        3.82,
        1.25,
        fill=WHITE,
        line=RGBColor(0xD2, 0xD0, 0xCE),
        line_width=1,
    )
    dot = add_box(
        slide,
        left + 0.2,
        5.32,
        0.12,
        0.12,
        fill=PURPLE,
        line=PURPLE,
        radius=False,
        line_width=0,
    )
    dot.line.fill.background()
    add_text(
        slide,
        title,
        left + 0.45,
        5.16,
        3.1,
        0.3,
        size=14,
        bold=True,
    )
    add_text(
        slide,
        body,
        left + 0.45,
        5.55,
        3.1,
        0.5,
        size=10.5,
        color=GRAY,
    )


def configure_enterprise_system_slide(slide):
    clear_slide(slide)
    background = add_box(
        slide,
        0,
        0,
        13.34,
        7.5,
        fill=WHITE,
        line=WHITE,
        radius=False,
        line_width=0,
    )
    background.line.fill.background()
    add_brand(slide)
    add_text(
        slide,
        "One enterprise agent system",
        0.62,
        0.72,
        8.7,
        0.55,
        size=30,
        bold=True,
    )
    add_text(
        slide,
        "Build \u2192 context \u2192 run \u2192 govern \u2192 improve \u2192 surface",
        0.62,
        1.35,
        10.5,
        0.36,
        size=13,
        color=GRAY,
        font_name="Aptos Mono",
    )

    stages = (
        ("GitHub", "Code \u00b7 prompts\nskills \u00b7 tools"),
        ("Microsoft IQ", "Enterprise data\nand context"),
        ("Foundry", "Agent runtime\nand models"),
        ("Entra + Purview", "Identity \u00b7 policy\nsafety \u00b7 approval"),
        ("Evaluate +\noptimize", "Quality \u00b7 cost\nfeedback loops"),
        ("Apps + M365", "APIs \u00b7 Teams\nwhere work happens"),
    )
    for index, (title, body) in enumerate(stages, start=1):
        left = 0.62 + (index - 1) * 2.07
        add_stage_card(
            slide,
            index,
            title,
            body,
            left,
            emphasized=index == 3,
        )
        if index < len(stages):
            arrow = add_box(
                slide,
                left + 1.82,
                3.02,
                0.22,
                0.34,
                fill=PURPLE,
                line=PURPLE,
                radius=False,
                line_width=0,
            )
            arrow._element.spPr.prstGeom.set("prst", "chevron")
            arrow.line.fill.background()

    add_value_card(
        slide,
        "Compose, don\u2019t replace",
        "Keep your gateway and cloud.\nAdd Foundry\u2019s agent factory.",
        0.62,
    )
    add_value_card(
        slide,
        "Govern by design",
        "Identity \u00b7 policy \u00b7 safety \u00b7 approval\nsurround the runtime.",
        4.75,
    )
    add_value_card(
        slide,
        "Improve under evidence",
        "Evaluation and telemetry\nfeed controlled releases.",
        8.88,
    )
    add_footer(slide, 2)


def add_pattern_chip(slide, text, left, top, width, *, fill=WHITE):
    add_box(
        slide,
        left,
        top,
        width,
        0.42,
        fill=fill,
        line=PURPLE,
        line_width=1,
    )
    add_text(
        slide,
        text,
        left + 0.08,
        top + 0.02,
        width - 0.16,
        0.36,
        size=9.5,
        color=NAVY,
        bold=True,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )


def configure_catalog_map_slide(slide):
    clear_slide(slide)
    background = add_box(
        slide,
        0,
        0,
        13.34,
        7.5,
        fill=WHITE,
        line=WHITE,
        radius=False,
        line_width=0,
    )
    background.line.fill.background()
    add_brand(slide)
    add_text(
        slide,
        "Where this catalog fits",
        0.62,
        0.72,
        8.7,
        0.55,
        size=30,
        bold=True,
    )
    add_text(
        slide,
        "Build in GitHub. Run and improve in Foundry. Reach users through channels you already own.",
        0.62,
        1.35,
        11.4,
        0.36,
        size=13,
        color=GRAY,
    )

    add_box(
        slide,
        0.62,
        2.0,
        2.0,
        4.05,
        fill=NAVY,
        line=NAVY,
        line_width=0,
    )
    add_text(
        slide,
        "BUILD IN\nGITHUB",
        0.86,
        2.28,
        1.52,
        0.64,
        size=17,
        color=WHITE,
        bold=True,
    )
    add_text(
        slide,
        "Code \u00b7 prompts\nskills \u00b7 graphs\nevaluation assets",
        0.86,
        3.12,
        1.52,
        0.72,
        size=10.5,
        color=WHITE,
    )
    build_chips = (
        "04  Loop",
        "05A  Agents",
        "05B  Workflows",
        "15  Lifecycle",
    )
    for index, chip in enumerate(build_chips):
        add_pattern_chip(
            slide,
            chip,
            0.84,
            4.02 + index * 0.46,
            1.56,
            fill=LIGHT_PURPLE,
        )

    frame = add_box(
        slide,
        2.96,
        1.92,
        7.42,
        4.22,
        fill=LIGHT_GRAY,
        line=PURPLE,
        line_width=2,
    )
    frame.adjustments[0] = 0.03
    header = add_box(
        slide,
        2.96,
        1.92,
        7.42,
        0.58,
        fill=PURPLE,
        line=PURPLE,
        radius=False,
        line_width=0,
    )
    header.line.fill.background()
    add_text(
        slide,
        "RUN + IMPROVE IN FOUNDRY \u2014 THE AGENT FACTORY",
        3.2,
        2.05,
        6.94,
        0.3,
        size=14,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    rows = (
        ("CONTEXT", ("03  Microsoft IQ", "10  Memory", "12  Toolboxes")),
        ("RUN", ("02  Agent Service", "09  Cross-cloud")),
        ("GOVERN", ("01  AI Gateway", "08  AI Safety", "13  Approval")),
        (
            "IMPROVE",
            ("06  Observability", "07  Eval gate", "11  Cost", "14  Adaptation"),
        ),
    )
    for row_index, (label, chips) in enumerate(rows):
        top = 2.75 + row_index * 0.78
        add_text(
            slide,
            label,
            3.18,
            top + 0.08,
            0.92,
            0.3,
            size=10,
            color=PURPLE,
            bold=True,
            font_name="Aptos Mono",
        )
        chip_width = (
            1.22
            if len(chips) == 4
            else 1.65
            if len(chips) == 3
            else 2.45
        )
        for chip_index, chip in enumerate(chips):
            add_pattern_chip(
                slide,
                chip,
                4.16 + chip_index * (chip_width + 0.18),
                top,
                chip_width,
            )

    for left in (2.68, 10.42):
        arrow = add_box(
            slide,
            left,
            3.62,
            0.2,
            0.42,
            fill=PURPLE,
            line=PURPLE,
            radius=False,
            line_width=0,
        )
        arrow._element.spPr.prstGeom.set("prst", "chevron")
        arrow.line.fill.background()

    add_box(
        slide,
        10.66,
        2.0,
        2.06,
        4.05,
        fill=LIGHT_PURPLE,
        line=PURPLE,
        line_width=1.5,
    )
    add_text(
        slide,
        "SURFACE",
        10.9,
        2.3,
        1.58,
        0.32,
        size=17,
        color=PURPLE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "Apps\nAPIs\nTeams\nM365\nother channels",
        10.9,
        3.0,
        1.58,
        1.55,
        size=13,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "Use the channels\nyour users already own.",
        10.9,
        5.05,
        1.58,
        0.55,
        size=9.5,
        color=GRAY,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "Primary placement only \u2014 the patterns compose across the lifecycle.",
        3.14,
        6.3,
        7.1,
        0.3,
        size=10,
        color=GRAY,
        align=PP_ALIGN.CENTER,
        font_name="Aptos Mono",
    )
    add_footer(slide, 3)


def ensure_framing_slides(presentation):
    entry = next(
        slide
        for slide in presentation.slides
        if any(
            value in slide_text(slide)
            for value in (
                "Keep your gateway and cloud \u2014 add the agent factory",
                "One enterprise agent system",
            )
        )
    )
    configure_enterprise_system_slide(entry)
    catalog = next(
        (
            slide
            for slide in presentation.slides
            if "Where this catalog fits" in slide_text(slide)
        ),
        None,
    )
    if catalog is None:
        catalog = presentation.slides.add_slide(presentation.slide_layouts[0])
    configure_catalog_map_slide(catalog)
    return entry, catalog


def configure_new_pattern_slide(
    slide,
    *,
    number,
    title,
    quote,
    demo,
    nodes,
    arrow_labels,
    benefits,
):
    shapes = list(slide.shapes)
    set_text(shapes[5], deck_label(number))
    set_text(shapes[6], "PATTERN")
    set_text(shapes[7], title)
    set_font_size(shapes[7], 27)
    set_text(shapes[8], quote)
    set_font_size(shapes[8], 16)
    set_text(shapes[11], demo)
    set_font_size(shapes[11], 13)

    for index, text in zip((13, 15, 17, 19), nodes):
        set_text(shapes[index], text)
        set_font_size(shapes[index], 14)
    for index, text in zip((21, 23), arrow_labels):
        set_text(shapes[index], text)
        set_font_size(shapes[index], 10.5)
        shapes[index].text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        shapes[index].text_frame.margin_left = 0
        shapes[index].text_frame.margin_right = 0
        shapes[index].text_frame.margin_top = 0
        shapes[index].text_frame.margin_bottom = 0
        shapes[index].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    # The template's labels were wider than the inter-box gaps, so the next
    # node clipped their final characters. Keep the first label inside the
    # 1.05-inch horizontal gap and the second clear of the right-hand node.
    shapes[21].left = Inches(3.72)
    shapes[21].top = Inches(3.43)
    shapes[21].width = Inches(0.96)
    shapes[21].height = Inches(0.48)
    shapes[21].text_frame.paragraphs[0].line_spacing = 1.05
    shapes[23].left = Inches(7.95)
    shapes[23].top = Inches(3.32)
    shapes[23].width = Inches(1.55)
    shapes[23].height = Inches(0.34)
    for index, text in zip((28, 31, 34), benefits):
        set_text(shapes[index], text)
        set_font_size(shapes[index], 13)


def ensure_pattern_five_variants(presentation):
    patterns = pattern_slides(presentation)
    agent_slide = patterns.get("5A") or patterns.get(5)
    if agent_slide is None:
        raise ValueError("Could not find the existing Pattern 5 slide")
    number_shape = next(
        shape
        for shape in agent_slide.shapes
        if getattr(shape, "has_text_frame", False)
        and shape.top / 914400 < 1.5
        and shape.text in {"05", "05A"}
    )
    set_text(number_shape, "05A")

    patterns = pattern_slides(presentation)
    if "5B" not in patterns:
        workflow_slide = clone_shape_slide(presentation, patterns[1])
        configure_new_pattern_slide(
            workflow_slide,
            number="5B",
            title="Workflow orchestration (graph-based pipeline)",
            quote=(
                "\u201cAgents contribute judgment. The workflow owns routing, "
                "state and control.\u201d"
            ),
            demo=(
                "Mixed code + agent graph \u00b7 deterministic switch \u00b7 "
                "fail-closed default \u00b7 checkpoint resume."
            ),
            nodes=(
                "Validate request\ncode executor",
                "Classify risk\nagent executor",
                "Policy + switch\nexplicit edges",
                "Audit record\ncode executor",
            ),
            arrow_labels=("typed\nmessage", "deterministic\nroute"),
            benefits=(
                "WorkflowBuilder makes nodes, edges and terminal output inspectable",
                "Production / impact rules override model judgment and route exceptions",
                "Trusted superstep checkpoints resume to the same audit ID and decision",
            ),
        )


def ensure_new_pattern_slides(presentation):
    patterns = pattern_slides(presentation)
    template = patterns[1]
    specifications = (
        {
            "number": 13,
            "title": "Human approval for consequential tool actions",
            "quote": "\u201cPause the exact tool call. Approval controls intent; authorization still controls access.\u201d",
            "demo": "Read runs immediately \u00b7 reject creates zero effects \u00b7 approve + replay creates exactly one.",
            "nodes": (
                "Prompt agent\nread + proposed action",
                "Foundry approval\nexact name + arguments",
                "Operator decision\nseparate identity",
                "Change-control MCP\none-time nonce",
            ),
            "arrow_labels": ("approval\nrequest", "approve / reject"),
            "benefits": (
                "Read-only tools continue without an approval interruption",
                "Rejected and stale decisions fail closed with zero side effects",
                "Nonce + decision + effect IDs correlate; replay stays exactly once",
            ),
        },
        {
            "number": 14,
            "title": "Model adaptation (fine-tuning & evaluation)",
            "quote": "\u201cFine-tune stable behavior, not changing knowledge \u2014 and prove the gain on untouched data.\u201d",
            "demo": "Base benchmark \u2192 reviewed SFT job \u2192 identical held-out eval \u2192 release gate + cleanup.",
            "nodes": (
                "Held-out test\nnever used for training",
                "Base model\nbenchmark first",
                "Foundry SFT job\nreviewed JSONL",
                "Developer eval tier\nsame test + cleanup",
            ),
            "arrow_labels": ("baseline\nmetrics", "stable behavior"),
            "benefits": (
                "RAG / Search / Foundry IQ remain the path for changing knowledge",
                "Schema, accuracy, adherence, tokens and latency \u2014 not training loss",
                "No measured gain or any configured regression means no promotion",
            ),
        },
        {
            "number": 15,
            "title": "Agent lifecycle & promotion (dev \u2192 test \u2192 prod)",
            "quote": "\u201cPromote immutable versions behind one endpoint; roll back the selector, not the state.\u201d",
            "demo": "Failing candidate blocked \u00b7 passing candidate pinned \u00b7 same URL \u00b7 rollback restores v1 + state.",
            "nodes": (
                "Release manifest\ncommit + aliases",
                "Dev \u2192 test\nsmoke + cloud eval",
                "Stable endpoint\nselector \u2192 candidate",
                "Rollback\nselector \u2192 prior",
            ),
            "arrow_labels": ("immutable\nversion", "eval evidence"),
            "benefits": (
                "Current agent object model \u2014 no new legacy Agent Application",
                "Stable endpoint URL does not change across promotion or rollback",
                "OIDC + complete release record; conversation store is not deleted",
            ),
        },
    )
    for specification in specifications:
        slide = patterns.get(specification["number"])
        if slide is None:
            slide = clone_shape_slide(presentation, template)
        configure_new_pattern_slide(slide, **specification)


def set_table_rows(shape, rows):
    table = shape.table
    xml_rows = list(table._tbl.tr_lst)
    while len(xml_rows) > len(rows):
        table._tbl.remove(xml_rows.pop())
    while len(xml_rows) < len(rows):
        clone = deepcopy(xml_rows[-1])
        table._tbl.append(clone)
        xml_rows.append(clone)

    for row_index, values in enumerate(rows):
        table.rows[row_index].height = Inches(0.46 if row_index else 0.42)
        for column_index, value in enumerate(values):
            set_text(table.cell(row_index, column_index), value)


def equalize_table_height(shape, inches):
    height = Inches(inches)
    row_height = int(height / len(shape.table.rows))
    for row in shape.table.rows:
        row.height = row_height
    shape.height = height


def merge_table_row(shape, row_index):
    table = shape.table
    first = table.cell(row_index, 0)
    if not first.is_merge_origin:
        first.merge(table.cell(row_index, 1))


def insert_table_row(shape, source_index, before_index):
    """Insert an unmerged row by cloning an existing data row."""
    table = shape.table
    source = list(table._tbl.tr_lst)[source_index]
    target = list(table._tbl.tr_lst)[before_index]
    target.addprevious(deepcopy(source))


def replace_table_row(shape, target_index, source_index):
    """Replace a previously merged row with a clean clone of an unmerged data row."""
    rows = list(shape.table._tbl.tr_lst)
    target = rows[target_index]
    target.addprevious(deepcopy(rows[source_index]))
    shape.table._tbl.remove(target)


def set_group_slide(slide, title, promise, pattern_list):
    texts = shape_texts(slide)
    old_title = next(
        value
        for value in texts
        if value
        in {
            "Control plane",
            "Agent factory",
            "Orchestration & interop",
            "Operate & optimise",
            "Platform foundation & governance",
            "Agent construction & knowledge",
            "Orchestration & interoperability",
            "Lifecycle, assurance & operations",
        }
    )
    replace_text(slide, old_title, title)
    candidates = [
        shape
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
        and 3.4 < shape.top / 914400 < 4.2
        and shape.text != title
    ]
    set_text(candidates[0], promise)
    list_shape = next(
        shape
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and 4.2 < shape.top / 914400 < 4.8
    )
    set_text(list_shape, pattern_list)
    set_font_size(list_shape, 13 if len(pattern_list) > 95 else 15)


def style_search_card(background, text_shape, live):
    if live:
        background.fill.solid()
        background.fill.fore_color.rgb = PURPLE
        background.line.color.rgb = PURPLE
        background.line.dash_style = MSO_LINE_DASH_STYLE.SOLID
        color = WHITE
    else:
        background.fill.solid()
        background.fill.fore_color.rgb = WHITE
        background.line.color.rgb = PURPLE
        background.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        color = PURPLE
    for paragraph in text_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = color


def update_pattern_three(slide):
    replace_text(
        slide,
        (
            "Web IQ published as our own MCP API on APIM \u2014 the gateway authenticates "
            "the caller, holds the key and meters every tool call.",
            "Web IQ via APIM + Azure AI Search via a Foundry agent \u2014 both live, both cited.",
        ),
        "Web IQ via APIM + Azure AI Search via a Foundry agent \u2014 both live, both cited.",
    )
    replace_text(
        slide,
        (
            "Web IQ (MCP) \u2014 cited web context",
            "Web IQ \u2014 LIVE via APIM",
        ),
        "Web IQ \u2014 LIVE via APIM",
    )
    replace_text(
        slide,
        (
            "Foundry IQ \u2014 Azure AI Search",
            "Azure AI Search \u2014 LIVE enterprise index",
        ),
        "Azure AI Search \u2014 LIVE enterprise index",
    )
    replace_text(
        slide,
        (
            "Work IQ \u2014 M365 org context",
            "Foundry IQ \u2014 managed KB (NARRATED)",
        ),
        "Foundry IQ \u2014 managed KB (NARRATED)",
    )
    replace_text(
        slide,
        (
            "Fabric IQ \u2014 Business context",
            "Fabric IQ + Work IQ \u2014 NARRATED",
        ),
        "Fabric IQ + Work IQ \u2014 NARRATED",
    )
    replace_text(
        slide,
        (
            "Caller sends an Entra token\nAPIM holds the Web IQ key",
            "APIM subscription key\nupstream Web IQ key stays on gateway",
        ),
        "APIM subscription key\nupstream Web IQ key stays on gateway",
    )
    replace_text(slide, ("no API key", "Basic v2 exception"), "Basic v2 exception")
    replace_text(
        slide,
        (
            "Governed MCP grounding \u2014 one AI-Gateway endpoint, any agent or model",
            "Web IQ on APIM \u2014 authenticated, metered, upstream key retained",
        ),
        "Web IQ on APIM \u2014 authenticated, metered, upstream key retained",
    )
    replace_text(
        slide,
        (
            "The key stays on APIM; callers authenticate with Entra and get metered",
            "AzureAISearchTool \u2014 project connection + cited enterprise index",
        ),
        "AzureAISearchTool \u2014 project connection + cited enterprise index",
    )
    replace_text(
        slide,
        (
            "All four IQ layers \u2014 world, enterprise, org and business meaning",
            "Foundry IQ, Fabric IQ and Work IQ \u2014 broader narrated layers",
        ),
        "Foundry IQ, Fabric IQ and Work IQ \u2014 broader narrated layers",
    )

    cards = sorted(
        [
            shape
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
            and shape.text
            and shape.left / 914400 > 8
            and 2.5 < shape.top / 914400 < 5.2
        ],
        key=lambda shape: shape.top,
    )
    backgrounds = []
    for card in cards:
        background = next(
            shape
            for shape in slide.shapes
            if not getattr(shape, "text", "")
            and abs(shape.left - (card.left - Inches(0.05))) < Inches(0.02)
            and abs(shape.top - card.top) < Inches(0.02)
        )
        backgrounds.append(background)
    for index, (background, card) in enumerate(zip(backgrounds, cards)):
        style_search_card(background, card, live=index < 2)

    connectors = sorted(
        [
            shape
            for shape in slide.shapes
            if not getattr(shape, "text", "")
            and 7.0 < shape.left / 914400 < 7.2
            and 2.5 < shape.top / 914400 < 5.2
        ],
        key=lambda shape: shape.top,
    )
    for index, connector in enumerate(connectors):
        connector.line.color.rgb = PURPLE
        connector.line.dash_style = (
            MSO_LINE_DASH_STYLE.SOLID if index < 2 else MSO_LINE_DASH_STYLE.DASH
        )

    if not any("SOLID = LIVE" in value for value in shape_texts(slide)):
        legend = slide.shapes.add_textbox(
            Inches(8.05),
            Inches(2.34),
            Inches(4.35),
            Inches(0.22),
        )
        legend.text_frame.margin_left = 0
        legend.text_frame.margin_right = 0
        legend.text_frame.margin_top = 0
        legend.text_frame.margin_bottom = 0
        legend.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        set_text(legend, "SOLID = LIVE    \u00b7    DASHED = NARRATED")
        for paragraph in legend.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.name = "Aptos"
                run.font.size = Pt(8)
                run.font.bold = True
                run.font.color.rgb = PURPLE


def update_pattern_five(slide):
    replace_text(
        slide,
        (
            "Multi-agent orchestration (Agent Framework)",
            "Agent orchestration (multi-agent coordination)",
        ),
        "Agent orchestration (multi-agent coordination)",
    )
    replace_text(
        slide,
        (
            "\u201cOrchestrate specialists when the problem is genuinely parallel \u2014 not because a framework let you.\u201d",
            "\u201cCoordinate agents with a reusable pattern; use a graph when the process must own control.\u201d",
        ),
        "\u201cCoordinate agents with a reusable pattern; use a graph when the process must own control.\u201d",
    )
    replace_text(
        slide,
        (
            "Request fans out to two specialists concurrently; Compliance returns BLOCK, then aggregate.",
            "Concurrent specialists, published as one Foundry-managed Responses endpoint.",
            "ConcurrentBuilder fans out to specialists; one hosted endpoint returns the fan-in result.",
        ),
        "ConcurrentBuilder fans out to specialists; one hosted endpoint returns the fan-in result.",
    )
    replace_text(
        slide,
        ("Orchestrator", "Foundry-hosted endpoint"),
        "Foundry-hosted endpoint",
    )
    replace_text(slide, ("Aggregate", "Fan-in result"), "Fan-in result")
    replace_text(
        slide,
        (
            "Open, code-first (SK / AutoGen lineage), model-portable",
            "Managed endpoint \u00b7 compute \u00b7 versions \u00b7 Entra Agent ID",
        ),
        "Managed endpoint \u00b7 compute \u00b7 versions \u00b7 Entra Agent ID",
    )
    replace_text(
        slide,
        (
            "Managed orchestration (concurrent / sequential / handoff / Magentic)",
            "ConcurrentBuilder \u00b7 high-level multi-agent collaboration pattern",
        ),
        "ConcurrentBuilder \u00b7 high-level multi-agent collaboration pattern",
    )
    replace_text(
        slide,
        (
            "Honest guidance: default to Pattern 4; escalate when earned",
            "5A coordinates agents \u00b7 5B controls an explicit process graph",
        ),
        "5A coordinates agents \u00b7 5B controls an explicit process graph",
    )


def update_content(presentation):
    patterns = pattern_slides(presentation)
    groups = group_slides(presentation)

    title_slide = presentation.slides[0]
    replace_text(
        title_slide,
        (
            "12 patterns   \u00b7   12 demos   \u00b7   for architects & principal engineers",
            "15 patterns   \u00b7   15 demos   \u00b7   for architects & principal engineers",
            "15 pattern families   \u00b7   16 demos   \u00b7   for architects & principal engineers",
        ),
        "15 pattern families   \u00b7   16 demos   \u00b7   for architects & principal engineers",
    )

    run_show = find_slide(presentation, "Run-of-show")
    tables = [shape for shape in run_show.shapes if getattr(shape, "has_table", False)]
    if len(tables[0].table.rows) == 9:
        insert_table_row(tables[0], source_index=3, before_index=4)
    # The source right-hand table has a merged section row at index 4. Insert a
    # normal data row before it so Pattern 9 keeps two independent cells.
    if len(tables[1].table.rows) == 8:
        insert_table_row(tables[1], source_index=3, before_index=4)
    if (
        len(tables[1].table.rows) >= 7
        and tables[1].table.cell(5, 0).is_merge_origin
    ):
        replace_table_row(tables[1], target_index=5, source_index=3)
    set_table_rows(
        tables[0],
        [
            ("Index", "Patterns"),
            ("PLATFORM FOUNDATION & GOVERNANCE", ""),
            ("1", "AI gateway & model access (APIM)"),
            ("8", "AI safety (Prompt Shields + Content Safety)"),
            ("13", "Human approval for consequential tool actions"),
            ("AGENT CONSTRUCTION & KNOWLEDGE", ""),
            ("2", "Foundry Agent Service (prompt and hosted agents)"),
            ("3", "Microsoft IQ \u2014 the intelligence layer"),
            ("12", "Centralized Toolboxes (one governed MCP endpoint)"),
            ("14", "Model adaptation (fine-tuning & evaluation)"),
            ("10", "Memory (short-term + long-term)"),
        ],
    )
    merge_table_row(tables[0], 1)
    merge_table_row(tables[0], 5)
    set_table_rows(
        tables[1],
        [
            ("Index", "Patterns"),
            ("ORCHESTRATION & INTEROPERABILITY", ""),
            ("4", "Agentic Loop (build skills, not agents)"),
            ("5A", "Agent orchestration (multi-agent coordination)"),
            ("5B", "Workflow orchestration (graph-based pipeline)"),
            ("9", "Cross-cloud interop (MCP / A2A)"),
            ("LIFECYCLE, ASSURANCE & OPERATIONS", ""),
            ("7", "Evaluation & release gate"),
            ("6", "Observability & tracing (OpenTelemetry)"),
            ("11", "Cost & latency (prompt cache + Model Router)"),
            ("15", "Agent lifecycle & promotion (dev \u2192 test \u2192 prod)"),
        ],
    )
    merge_table_row(tables[1], 1)
    merge_table_row(tables[1], 6)
    equalize_table_height(tables[0], 4.48)
    equalize_table_height(tables[1], 4.48)
    replace_text(
        run_show,
        (
            "Twelve patterns in four groups. One slide, one live demo each.",
            "Fifteen patterns in four groups. One slide, one live demo each.",
            "Fifteen pattern families in four groups. Sixteen runnable demos.",
        ),
        "Fifteen pattern families in four groups. Sixteen runnable demos.",
    )
    close = find_slide(presentation, "and the close")
    close_message = (
        "Your other cloud runs an agent. Foundry runs the agent FACTORY. "
        "MCP + A2A join the two.\n"
        "Keep your existing stack where it works \u00b7 add Foundry where you have gaps "
        "(identity, eval, tracing, grounding, safety) \u00b7 unify governance with "
        "Entra and Purview across your estate."
    )
    close_body = next(
        shape
        for shape in close.shapes
        if getattr(shape, "has_text_frame", False)
        and shape.text.startswith("Your other cloud runs an agent.")
    )
    set_text(close_body, close_message)
    for index, paragraph in enumerate(close_body.text_frame.paragraphs):
        for run in paragraph.runs:
            run.font.size = Pt(21 if index == 0 else 15)

    set_group_slide(
        groups[1],
        "Platform foundation & governance",
        "Establish governed model access and safety controls",
        "01  AI gateway   \u00b7   08  AI safety   \u00b7   13  Human approval",
    )
    set_group_slide(
        groups[2],
        "Agent construction & knowledge",
        "Build agents with governed knowledge, tools and memory",
        "02  Agent Service  \u00b7  03  Microsoft IQ  \u00b7  12  Toolboxes  \u00b7  14  Adaptation  \u00b7  10  Memory",
    )
    set_group_slide(
        groups[3],
        "Orchestration & interoperability",
        "Compose skills, specialists and cross-cloud systems",
        "04  Agentic Loop   \u00b7   05A  Agent orchestration   \u00b7   05B  Workflow orchestration   \u00b7   09  Cross-cloud",
    )
    set_group_slide(
        groups[4],
        "Lifecycle, assurance & operations",
        "Evaluate, observe and optimize every release",
        "07  Evaluation   \u00b7   06  Observability   \u00b7   11  Cost   \u00b7   15  Lifecycle",
    )

    replace_text(
        patterns[1],
        ("Wedge \u2192 AI Hub Gateway / Citadel", "AI gateway & model access (APIM)"),
        "AI gateway & model access (APIM)",
    )
    replace_text(
        patterns[1],
        (
            "Ships as AI Hub Gateway / Citadel (Foundry + APIM)",
            "Enterprise APIM + Foundry composition \u2014 no rip-and-replace",
        ),
        "Enterprise APIM + Foundry composition \u2014 no rip-and-replace",
    )
    replace_text(
        patterns[8],
        (
            "Governance (Prompt Shields + Content Safety)",
            "AI safety (Prompt Shields + Content Safety)",
        ),
        "AI safety (Prompt Shields + Content Safety)",
    )
    replace_text(
        patterns[2],
        (
            "Agent Service (prompt and hosted agent)",
            "Foundry Agent Service (prompt and hosted agents)",
        ),
        "Foundry Agent Service (prompt and hosted agents)",
    )
    update_pattern_three(patterns[3])
    update_pattern_five(patterns["5A"])
    replace_text(
        patterns[7],
        ("Evaluation \u2192 optimization (CI gate)", "Evaluation & release gate"),
        "Evaluation & release gate",
    )
    replace_text(
        patterns[7],
        (
            "Score the golden set locally AND in Foundry; a wrong 'suitable' tanks "
            "groundedness and the CI gate blocks the PR.",
            "Foundry cloud scorecard + CI release gate; a wrong answer blocks the PR.",
            "Generate candidate answers, score in Foundry, fail closed; demo mode plants one regression.",
        ),
        "Generate candidate answers, score in Foundry, fail closed; demo mode plants one regression.",
    )
    replace_text(
        patterns[7],
        (
            "Golden set\n+ planted wrong row",
            "Golden set\nquestions + policy context",
        ),
        "Golden set\nquestions + policy context",
    )
    replace_text(
        patterns[7],
        ("Agent", "Candidate\nprompt + model"),
        "Candidate\nprompt + model",
    )
    replace_text(
        patterns[7],
        (
            "Evaluators\ngroundedness \u00b7 tool-accuracy",
            "Evaluators\ngroundedness \u00b7 relevance \u00b7 coherence",
        ),
        "Evaluators\ngroundedness \u00b7 relevance \u00b7 coherence",
    )
    replace_text(
        patterns[7],
        (
            "Scores locally + uploads to the Foundry Evaluations tab",
            "Candidate answers are generated before cloud evaluation",
        ),
        "Candidate answers are generated before cloud evaluation",
    )
    replace_text(
        patterns[7],
        (
            "CI gate on every PR \u2014 regressions can't merge",
            "Scoped CI gate \u2014 relevant regressions can't merge",
        ),
        "Scoped CI gate \u2014 relevant regressions can't merge",
    )
    replace_text(
        patterns[6],
        (
            "Run enable_tracing.py \u2014 agent 'rm-assistant-traced' + its trace in BOTH "
            "Foundry Tracing and App Insights.",
            "Metadata-only OTel trace in Foundry + App Insights; content requires explicit opt-in.",
        ),
        "Metadata-only OTel trace in Foundry + App Insights; content requires explicit opt-in.",
    )
    replace_text(
        patterns[6],
        (
            "OpenTelemetry spans\ntokens \u00b7 latency \u00b7 cost",
            "OpenTelemetry spans\nmetadata \u00b7 tokens \u00b7 latency",
            "OpenTelemetry spans\nmetadata only \u00b7 tokens \u00b7 latency",
        ),
        "OpenTelemetry spans\nmetadata only \u00b7 tokens \u00b7 latency",
    )
    replace_text(
        patterns[6],
        (
            "Token + cost + latency per span \u2014 observability & FinOps",
            "Token + latency telemetry \u00b7 rate-card cost by agent/version",
        ),
        "Token + latency telemetry \u00b7 rate-card cost by agent/version",
    )
    outcome_label = next(
        shape
        for shape in patterns[6].shapes
        if getattr(shape, "has_text_frame", False)
        and shape.text == "Foundry \u2014 Agents + Tracing"
    )
    if outcome_label.top / 914400 > 4.7:
        for shape in patterns[6].shapes:
            if 4.5 < shape.top / 914400 < 5.5:
                shape.top -= Inches(0.22)
    dashed_card = next(
        shape
        for shape in patterns[6].shapes
        if not getattr(shape, "text", "")
        and shape.left / 914400 > 7
        and shape.top / 914400 > 4
        and shape.line.dash_style == MSO_LINE_DASH_STYLE.DASH
        and shape.element.spPr.prstGeom.get("prst") == "roundRect"
    )
    dashed_card.left = Inches(8.99)
    dashed_card.top = Inches(4.56)
    dashed_card.width = Inches(3.6)
    dashed_card.height = Inches(0.7)
    dashed_connector = next(
        shape
        for shape in patterns[6].shapes
        if not getattr(shape, "text", "")
        and shape.line.dash_style == MSO_LINE_DASH_STYLE.DASH
        and shape.element.spPr.prstGeom.get("prst") == "line"
    )
    dashed_connector.left = Inches(9.0)
    dashed_connector.top = Inches(4.56)
    dashed_connector.width = Inches(1.79)
    dashed_connector.height = 0
    set_font_size(
        next(
            shape
            for shape in groups[4].shapes
            if getattr(shape, "has_text_frame", False)
            and shape.text == "Lifecycle, assurance & operations"
        ),
        39,
    )


def reorder_slides(presentation):
    patterns = pattern_slides(presentation)
    groups = group_slides(presentation)
    title = presentation.slides[0]
    entry = find_slide(presentation, "One enterprise agent system")
    catalog_map = find_slide(presentation, "Where this catalog fits")
    run_show = find_slide(presentation, "Run-of-show")
    close = find_slide(presentation, "and the close")
    order = [
        title,
        entry,
        catalog_map,
        run_show,
        groups[1],
        patterns[1],
        patterns[8],
        patterns[13],
        groups[2],
        patterns[2],
        patterns[3],
        patterns[12],
        patterns[14],
        patterns[10],
        groups[3],
        patterns[4],
        patterns["5A"],
        patterns["5B"],
        patterns[9],
        groups[4],
        patterns[7],
        patterns[6],
        patterns[11],
        patterns[15],
        close,
    ]
    slide_ids = [slide.slide_id for slide in order]
    id_to_element = {
        int(element.get("id")): element for element in presentation.slides._sldIdLst
    }
    for element in list(presentation.slides._sldIdLst):
        presentation.slides._sldIdLst.remove(element)
    for slide_id in slide_ids:
        presentation.slides._sldIdLst.append(id_to_element[slide_id])


def update_footers(presentation):
    for slide_number, slide in enumerate(presentation.slides, start=1):
        candidates = [
            shape
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
            and shape.left / 914400 > 11.5
            and shape.top / 914400 > 6.8
            and shape.text.strip().isdigit()
        ]
        for shape in candidates:
            set_text(shape, str(slide_number))


def main():
    presentation = Presentation(DECK)
    ensure_framing_slides(presentation)
    ensure_pattern_five_variants(presentation)
    ensure_new_pattern_slides(presentation)
    update_content(presentation)
    reorder_slides(presentation)
    update_footers(presentation)
    presentation.save(DECK)
    print(
        f"Refreshed {DECK.name}: 25 slides, "
        "15 pattern families and 16 demos."
    )


if __name__ == "__main__":
    main()
